"""
PowerPoint-Renderer — interner Microservice, der das von Straton produzierte Slide-HTML
(`<div data-theme="…"><section class="slide" data-layout="…">…</section>…</div>`, siehe
`src/features/chat/constants/pptxExportPrompt.ts`) in eine echte `.pptx`-Datei umwandelt.

Wird ausschliesslich von der Edge Function `generate-pptx-from-outline` aufgerufen (Shared-Secret
im Header `X-Internal-Token`) — kein öffentlich erreichbarer Endpoint.
"""

import math
import os
from io import BytesIO

from bs4 import BeautifulSoup
from docx import Document
from docx.enum.section import WD_SECTION_START
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Pt as DocxPt, RGBColor as DocxRGB
from fastapi import FastAPI, Header, HTTPException, Response
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pydantic import BaseModel

INTERNAL_TOKEN = os.environ.get("PPTX_RENDER_SERVICE_TOKEN", "")

ALLOWED_LAYOUTS = {"title", "section", "content", "table", "stats", "twocol", "agenda", "boxes"}
MAX_SLIDES = 60

# Pendant zu `PPTX_ICON_WHITELIST` in `pptxExportPrompt.ts` — gleiche Zeichen, gleiche Reihenfolge.
ICON_WHITELIST = {
    "🎯", "💡", "📈", "📊", "🔒", "🌍", "🚀", "⚡",
    "🤝", "💰", "✅", "⭐", "🛡", "🧩", "🔄", "📌",
    "⏱", "🧠", "🌱", "🏆", "🔧", "📍", "🔥", "🎓",
}


def validate_icon(text: str) -> str:
    """Nur whitelisted Icons rendern — alles andere wird stillschweigend zu leerem String (kein Crash)."""
    cleaned = text.strip()
    return cleaned if cleaned in ICON_WHITELIST else ""

# 96px = 1in in der Vorschau-CSS (`PPTX_SLIDE_NATIVE_WIDTH/HEIGHT`) — exakt 96 DPI, daher
# rechnen alle Mass-/Schriftangaben hier 1:1 aus den px-Werten in pptxOutline.ts um.
SLIDE_WIDTH_IN = 1280 / 96
SLIDE_HEIGHT_IN = 720 / 96
MARGIN_X_IN = 96 / 96
MARGIN_Y_IN = 72 / 96
CONTENT_WIDTH_IN = SLIDE_WIDTH_IN - 2 * MARGIN_X_IN
CONTENT_HEIGHT_IN = SLIDE_HEIGHT_IN - 2 * MARGIN_Y_IN

COLOR_TEXT_DARK = RGBColor(0x0F, 0x17, 0x2A)
COLOR_TEXT_LIGHT = RGBColor(0xF8, 0xFA, 0xFC)
COLOR_SUBTITLE_LIGHT = RGBColor(0xE2, 0xE8, 0xF0)
COLOR_BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_STAT_LABEL = RGBColor(0x47, 0x55, 0x69)

FONT_FAMILY = "Calibri"

# px → pt: 96px = 1in = 72pt, also px * 0.75.
PT_H1 = Pt(56 * 0.75)
PT_H2 = Pt(36 * 0.75)
PT_SUBTITLE = Pt(24 * 0.75)
PT_BODY = Pt(26 * 0.75)
PT_TABLE = Pt(22 * 0.75)
PT_COL_H2 = Pt(27)
PT_STAT_VALUE = Pt(40)
PT_STAT_LABEL = Pt(15)
PT_BOX_ICON = Pt(32)
PT_BOX_TITLE = Pt(18)
PT_BOX_TEXT = Pt(14)

DEFAULT_THEME = "blue"

# 5 kuratierte Paletten — Pendant zu `PPTX_THEME_PALETTES` in `pptxOutline.ts` (gleiche Hex-Werte).
THEME_PALETTES = {
    "blue": {
        "accent": RGBColor(0x1D, 0x4E, 0xD8),
        "accent_on_dark": RGBColor(0xBF, 0xDB, 0xFE),
        "bg_dark_from": RGBColor(0x0F, 0x17, 0x2A),
        "bg_dark_to": RGBColor(0x1E, 0x3A, 0x8A),
        "box_colors": (RGBColor(0x1D, 0x4E, 0xD8), RGBColor(0x0E, 0xA5, 0xE9), RGBColor(0x4F, 0x46, 0xE5)),
    },
    "green": {
        "accent": RGBColor(0x15, 0x80, 0x3D),
        "accent_on_dark": RGBColor(0xBB, 0xF7, 0xD0),
        "bg_dark_from": RGBColor(0x05, 0x2E, 0x1F),
        "bg_dark_to": RGBColor(0x15, 0x80, 0x3D),
        "box_colors": (RGBColor(0x15, 0x80, 0x3D), RGBColor(0x0D, 0x94, 0x88), RGBColor(0x65, 0xA3, 0x0D)),
    },
    "violet": {
        "accent": RGBColor(0x7C, 0x3A, 0xED),
        "accent_on_dark": RGBColor(0xDD, 0xD6, 0xFE),
        "bg_dark_from": RGBColor(0x2E, 0x10, 0x65),
        "bg_dark_to": RGBColor(0x6D, 0x28, 0xD9),
        "box_colors": (RGBColor(0x7C, 0x3A, 0xED), RGBColor(0xC0, 0x26, 0xD3), RGBColor(0x43, 0x38, 0xCA)),
    },
    "orange": {
        "accent": RGBColor(0xC2, 0x41, 0x0C),
        "accent_on_dark": RGBColor(0xFE, 0xD7, 0xAA),
        "bg_dark_from": RGBColor(0x43, 0x14, 0x07),
        "bg_dark_to": RGBColor(0xC2, 0x41, 0x0C),
        "box_colors": (RGBColor(0xC2, 0x41, 0x0C), RGBColor(0xD9, 0x77, 0x06), RGBColor(0xBE, 0x12, 0x3C)),
    },
    "slate": {
        "accent": RGBColor(0x33, 0x41, 0x55),
        "accent_on_dark": RGBColor(0xCB, 0xD5, 0xE1),
        "bg_dark_from": RGBColor(0x0F, 0x17, 0x2A),
        "bg_dark_to": RGBColor(0x33, 0x41, 0x55),
        "box_colors": (RGBColor(0x33, 0x41, 0x55), RGBColor(0x0F, 0x76, 0x6E), RGBColor(0xB4, 0x53, 0x09)),
    },
    "red": {
        "accent": RGBColor(0xB9, 0x1C, 0x1C),
        "accent_on_dark": RGBColor(0xFE, 0xCA, 0xCA),
        "bg_dark_from": RGBColor(0x45, 0x0A, 0x0A),
        "bg_dark_to": RGBColor(0xB9, 0x1C, 0x1C),
        "box_colors": (RGBColor(0xB9, 0x1C, 0x1C), RGBColor(0xEA, 0x58, 0x0C), RGBColor(0xDB, 0x27, 0x77)),
    },
    "pink": {
        "accent": RGBColor(0xBE, 0x18, 0x5D),
        "accent_on_dark": RGBColor(0xFB, 0xCF, 0xE8),
        "bg_dark_from": RGBColor(0x50, 0x07, 0x24),
        "bg_dark_to": RGBColor(0xBE, 0x18, 0x5D),
        "box_colors": (RGBColor(0xBE, 0x18, 0x5D), RGBColor(0xC0, 0x26, 0xD3), RGBColor(0xE1, 0x1D, 0x48)),
    },
    "teal": {
        "accent": RGBColor(0x0F, 0x76, 0x6E),
        "accent_on_dark": RGBColor(0x99, 0xF6, 0xE4),
        "bg_dark_from": RGBColor(0x04, 0x2F, 0x2E),
        "bg_dark_to": RGBColor(0x0F, 0x76, 0x6E),
        "box_colors": (RGBColor(0x0F, 0x76, 0x6E), RGBColor(0x08, 0x91, 0xB2), RGBColor(0x05, 0x96, 0x69)),
    },
    "amber": {
        "accent": RGBColor(0xB4, 0x53, 0x09),
        "accent_on_dark": RGBColor(0xFD, 0xE6, 0x8A),
        "bg_dark_from": RGBColor(0x45, 0x1A, 0x03),
        "bg_dark_to": RGBColor(0xB4, 0x53, 0x09),
        "box_colors": (RGBColor(0xB4, 0x53, 0x09), RGBColor(0xCA, 0x8A, 0x04), RGBColor(0xC2, 0x41, 0x0C)),
    },
    "indigo": {
        "accent": RGBColor(0x43, 0x38, 0xCA),
        "accent_on_dark": RGBColor(0xC7, 0xD2, 0xFE),
        "bg_dark_from": RGBColor(0x1E, 0x1B, 0x4B),
        "bg_dark_to": RGBColor(0x43, 0x38, 0xCA),
        "box_colors": (RGBColor(0x43, 0x38, 0xCA), RGBColor(0x4F, 0x46, 0xE5), RGBColor(0x63, 0x66, 0xF1)),
    },
}


# Pendant zu `PPTX_PRESET_SPECS` in `pptxOutline.ts` (gleiche Hex-Werte) — der NUTZER wählt eines
# über das Preset-Modal vor der Generierung, die KI wählt hier nichts mehr selbst (anders als
# `THEME_PALETTES`, das weiterhin für die Element-Akzentfarbe `data-color`/`data-textcolor` frei
# von der KI gewählt wird, siehe `extract_text_style`/`extract_container_style` — unverändert).
PRESET_SPECS = {
    "tech": {
        # Werte 1:1 aus dem Referenzbild gepixelt (Logo-Punkt/Überschrift ≈ rgb(0,195,198),
        # Tropfen-Glanzlicht ≈ rgb(21,157,179), Hintergrund ≈ rgb(6,12,28)) — siehe `decoration_style`.
        "accent": RGBColor(0x06, 0xC2, 0xC2),
        "accent_on_dark": RGBColor(0x7E, 0xEC, 0xEC),
        "bg_dark_from": RGBColor(0x05, 0x07, 0x0D),
        "bg_dark_to": RGBColor(0x0A, 0x16, 0x26),
        "box_colors": (RGBColor(0x06, 0xC2, 0xC2), RGBColor(0x08, 0x91, 0xB2), RGBColor(0x0E, 0x74, 0x90)),
        "heading_font": "Calibri",
        "corner_scale": 0.7,
        "title_treatment": "gradient-cover",
        "accent_spine": True,
        "decoration_style": "blob",
        # Ganzes Deck dunkel (nicht nur die Titelfolie) — heller Text, dunkel-transluzente Karten,
        # Teal-Akzente, Deko-Muster. Nur Tech; alle anderen Presets/Legacy bleiben hell (Feld fehlt).
        "surface": "dark",
    },
    "soft": {
        "accent": RGBColor(0xEC, 0x48, 0x99),
        "accent_on_dark": RGBColor(0xFB, 0xCF, 0xE8),
        "bg_dark_from": RGBColor(0xFD, 0xF2, 0xF8),
        "bg_dark_to": RGBColor(0xED, 0xE9, 0xFE),
        "box_colors": (RGBColor(0xEC, 0x48, 0x99), RGBColor(0xF4, 0x72, 0xB6), RGBColor(0xA7, 0x8B, 0xFA)),
        "heading_font": "Calibri",
        "corner_scale": 1.5,
        "title_treatment": "editorial-light",
        "accent_spine": False,
    },
    "professional": {
        "accent": RGBColor(0x1E, 0x3A, 0x5F),
        "accent_on_dark": RGBColor(0xCB, 0xD5, 0xE1),
        "bg_dark_from": RGBColor(0x0B, 0x12, 0x20),
        "bg_dark_to": RGBColor(0x1E, 0x3A, 0x5F),
        "box_colors": (RGBColor(0x1E, 0x3A, 0x5F), RGBColor(0x33, 0x41, 0x55), RGBColor(0x47, 0x55, 0x69)),
        "heading_font": "Calibri",
        "corner_scale": 0.7,
        "title_treatment": "gradient-cover",
        "accent_spine": True,
    },
    "bold": {
        "accent": RGBColor(0xF9, 0x73, 0x16),
        "accent_on_dark": RGBColor(0xFE, 0xD7, 0xAA),
        "bg_dark_from": RGBColor(0x1A, 0x0B, 0x2E),
        "bg_dark_to": RGBColor(0xBE, 0x18, 0x5D),
        "box_colors": (RGBColor(0xF9, 0x73, 0x16), RGBColor(0xDB, 0x27, 0x77), RGBColor(0x7C, 0x3A, 0xED)),
        "heading_font": "Calibri",
        "corner_scale": 1.15,
        "title_treatment": "gradient-cover",
        "accent_spine": True,
    },
    "minimal": {
        "accent": RGBColor(0x11, 0x18, 0x27),
        "accent_on_dark": RGBColor(0x9C, 0xA3, 0xAF),
        "bg_dark_from": RGBColor(0xF8, 0xFA, 0xFC),
        "bg_dark_to": RGBColor(0xF1, 0xF5, 0xF9),
        "box_colors": (RGBColor(0x11, 0x18, 0x27), RGBColor(0x37, 0x41, 0x51), RGBColor(0x6B, 0x72, 0x80)),
        "heading_font": "Georgia",
        "corner_scale": 0.4,
        "title_treatment": "editorial-light",
        "accent_spine": False,
    },
}

# Design-Spec für ein altes, reines `theme`-Deck — identisch zu den bisherigen Festwerten (Bestandsschutz).
LEGACY_THEME_DESIGN_DEFAULTS = {
    "heading_font": FONT_FAMILY,
    "corner_scale": 1.0,
    "title_treatment": "gradient-cover",
    "accent_spine": True,
}


def resolve_palette(theme_or_preset: str) -> dict:
    """Liest `preset` zuerst (neue Decks), fällt sonst auf das alte `theme`-System zurück (alte Decks, unverändert)."""
    if theme_or_preset in PRESET_SPECS:
        return PRESET_SPECS[theme_or_preset]
    base = THEME_PALETTES.get(theme_or_preset, THEME_PALETTES[DEFAULT_THEME])
    return {**base, **LEGACY_THEME_DESIGN_DEFAULTS}


def corner(palette: dict, default: float) -> float:
    """Skaliert einen Default-Eckenradius (Fraktion, kein `data-radius`-Override) mit dem Preset-`corner_scale`."""
    return default * palette.get("corner_scale", 1.0)


def blend_colors(rgb_a: RGBColor, rgb_b: RGBColor, fraction: float) -> RGBColor:
    """`fraction=1` → `rgb_a`, `fraction=0` → `rgb_b` — linearer Misch-Helfer für beide Richtungen."""
    r = round(rgb_a[0] * fraction + rgb_b[0] * (1 - fraction))
    g = round(rgb_a[1] * fraction + rgb_b[1] * (1 - fraction))
    b = round(rgb_a[2] * fraction + rgb_b[2] * (1 - fraction))
    return RGBColor(r, g, b)


def mix_with_white(rgb: RGBColor, accent_fraction: float) -> RGBColor:
    """Entspricht CSS `color-mix(in srgb, accent X%, white)` — X% Akzent, Rest weiss aufgehellt."""
    return blend_colors(rgb, COLOR_BG_WHITE, accent_fraction)


# --- Dunkles Deck (nur Tech-Preset, `surface == "dark"`) -------------------------------------------
# Pendant zur CSS-Verzweigung `isDark` in `buildPptxSlideThemeCss` (`pptxOutline.ts`). Alle Layouts
# (nicht nur die Titelfolie) rendern auf dunklem Grund: heller Text, dunkel-transluzente Karten,
# Teal-Akzente und Deko-Muster. Strikt auf `surface == "dark"` gegated, damit die anderen Presets
# (hell) byte-identisch wie bisher rendern.
COLOR_MUTED_ON_DARK = RGBColor(0x94, 0xA3, 0xB8)  # slate-400 — Label/Sekundärtext auf Dunkel


def is_dark_deck(palette: dict) -> bool:
    return palette.get("surface") == "dark"


def dark_card_fill(palette: dict, accent_fraction: float = 0.12) -> RGBColor:
    """Dunkel-transluzente Karte: kleiner Akzentanteil in den dunklen Hintergrund gemischt
    (Pendant zu CSS `color-mix(in srgb, accent X%, gradientTo)`)."""
    return blend_colors(palette["accent"], palette["bg_dark_to"], accent_fraction)


def add_card_hairline(slide, *, left_in: float, top_in: float, width_in: float, color: RGBColor) -> None:
    """Dünne Teal-Akzentkante oben auf einer dunklen Karte (statt Vollfarb-Block)."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left_in), Inches(top_in), Inches(width_in), Inches(0.045))
    bar.line.fill.background()
    bar.shadow.inherit = False
    bar.fill.solid()
    bar.fill.fore_color.rgb = color


# Deko-Platzierungen für das dunkle Tech-Deck — Pendant zu `PPTX_DARK_BAND_PLACEMENTS` in
# `pptxOutline.ts` (px → Inch umgerechnet, 96 dpi). Pro Inhalts-Folie wird EIN Dot-Band gezeichnet,
# dessen Position über den Folien-Index rotiert (`variant = index % len`), damit kein Muster
# „immer gleich" bleibt. `vertical` kippt das Band hochkant; `length` ist Breite (h) bzw. Höhe (v).
DOT_BAND_PLACEMENTS = [
    {"left": 8.27, "top": 6.15, "length": 4.48, "vertical": False},
    {"left": 9.23, "top": 0.40, "length": 3.44, "vertical": False},
    {"left": 12.15, "top": 1.56, "length": 3.75, "vertical": True},
    {"left": 0.58, "top": 6.19, "length": 3.75, "vertical": False},
    {"left": 3.96, "top": 6.27, "length": 4.48, "vertical": False},
    {"left": 0.42, "top": 1.56, "length": 3.54, "vertical": True},
]


def add_dot_band(slide, palette: dict, placement: dict, *,
                 dots: int = 26, rows: int = 3, amplitude_in: float = 0.26,
                 dot_in: float = 0.05, row_gap_in: float = 0.13) -> None:
    """Band aus Teal-Punkten entlang einer Sinuskurve — heller zur Bandmitte, an den Enden
    ausgeblendet (Pendant zum maskierten Dot-Band der Preview-CSS). Position/Orientierung kommen
    aus `placement` (siehe `DOT_BAND_PLACEMENTS`)."""
    vertical = placement.get("vertical", False)
    base_left = placement["left"]
    base_top = placement["top"]
    length = placement["length"]
    for row in range(rows):
        off = row * row_gap_in
        for i in range(dots):
            t = i / max(dots - 1, 1)
            edge = 1.0 - abs(0.5 - t) * 2.0
            frac = 0.2 + 0.55 * edge
            wave = math.sin(t * math.pi * 2.0) * amplitude_in
            if vertical:
                x = base_left + off + wave
                y = base_top + t * length
            else:
                x = base_left + t * length
                y = base_top + off + wave
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(dot_in), Inches(dot_in))
            dot.line.fill.background()
            dot.shadow.inherit = False
            dot.fill.solid()
            dot.fill.fore_color.rgb = blend_colors(palette["accent_on_dark"], palette["bg_dark_to"], frac)


def add_glow_ring(slide, palette: dict, *, cx_in: float, cy_in: float, radius_in: float) -> None:
    """Konzentrische Teal-Ringe hinter einer grossen Kennzahl (stats) — Deckkraft fällt nach aussen."""
    for k, frac in enumerate((0.8, 0.5, 0.28)):
        r = radius_in * (1.0 + k * 0.16)
        ring = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(cx_in - r), Inches(cy_in - r), Inches(2 * r), Inches(2 * r)
        )
        ring.fill.background()
        ring.shadow.inherit = False
        ring.line.color.rgb = blend_colors(palette["accent_on_dark"], palette["bg_dark_to"], frac)
        ring.line.width = Pt(2.0)


# Pendant zu `PPTX_RADIUS_KEYS`/`PPTX_SIZE_KEYS` in `pptxExportPrompt.ts` und den CSS-Regeln in
# `pptxOutline.ts` (`buildPptxElementStyleOverrideCss`) — gleiche Stufen, gleiche Faktoren.
RADIUS_ADJUSTMENTS = {"none": 0.0, "sm": 0.04, "md": 0.08, "lg": 0.14, "full": 0.5}
SIZE_FACTORS = {"sm": 0.75, "md": 1.0, "lg": 1.3, "xl": 1.6}
ALIGN_MAP = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}


def extract_text_style(el) -> dict:
    """Liest die optionalen `data-*`-Design-Attribute eines Text-Tags (siehe `ALLOWED_ELEMENT_STYLE_ATTRS`
    in `pptxOutline.ts`) — ungültige/fehlende Werte werden stillschweigend ignoriert (Standard-Look).
    `data-color`/`data-radius` sind nur für `h1` einer `title`-Folie relevant (Titel als farbige Box,
    siehe `render_title_slide`) — auf allen anderen Text-Tags werden sie schlicht nicht ausgewertet."""
    if el is None:
        return {}
    style: dict = {}
    factor = SIZE_FACTORS.get(el.get("data-size") or "")
    if factor is not None:
        style["size_factor"] = factor
    if el.get("data-bold") == "true":
        style["bold"] = True
    if el.get("data-italic") == "true":
        style["italic"] = True
    if el.get("data-underline") == "true":
        style["underline"] = True
    palette = THEME_PALETTES.get(el.get("data-textcolor") or "")
    if palette:
        style["color"] = palette["accent"]
    alignment = ALIGN_MAP.get(el.get("data-align") or "")
    if alignment is not None:
        style["alignment"] = alignment
    if (el.get("data-color") or "") in THEME_PALETTES:
        style["color_key"] = el.get("data-color")
    radius = el.get("data-radius")
    if radius in RADIUS_ADJUSTMENTS:
        style["radius"] = radius
    return style


def extract_container_style(el) -> dict:
    """Pendant zu `extract_text_style`, aber für Karten-/Gruppen-Container (`data-color`/`data-radius`/`data-valign`/`data-align`)."""
    if el is None:
        return {}
    style: dict = {}
    if (el.get("data-color") or "") in THEME_PALETTES:
        style["color_key"] = el.get("data-color")
    radius = RADIUS_ADJUSTMENTS.get(el.get("data-radius") or "")
    if radius is not None:
        style["radius"] = radius
    valign = el.get("data-valign")
    if valign in ("top", "middle", "bottom"):
        style["valign"] = valign
    alignment = ALIGN_MAP.get(el.get("data-align") or "")
    if alignment is not None:
        style["alignment"] = alignment
    return style


def parse_content_items(container) -> list[dict]:
    """`h1`/`h2`/`subtitle`/`p`/`ul`/`ol`/`table`/`callout` direkt unter `container` → typisierte Item-Liste."""
    items: list[dict] = []
    for child in container.find_all(
        ["h1", "h2", "subtitle", "p", "ul", "ol", "table", "callout"], recursive=False
    ):
        tag = child.name
        if tag == "callout":
            icon_el = child.find("icon", recursive=False)
            icon = validate_icon(icon_el.get_text(strip=True)) if icon_el else ""
            if icon_el:
                icon_el.extract()
            items.append(
                {
                    "type": tag,
                    "text": child.get_text(strip=True),
                    "icon": icon,
                    "container_style": extract_container_style(child),
                }
            )
        elif tag in ("h1", "h2", "subtitle", "p"):
            items.append({"type": tag, "text": child.get_text(strip=True), "style": extract_text_style(child)})
        elif tag in ("ul", "ol"):
            ordered = tag == "ol"
            for i, li in enumerate(child.find_all("li", recursive=False), start=1):
                items.append(
                    {
                        "type": "li",
                        "text": li.get_text(strip=True),
                        "ordered": ordered,
                        "index": i,
                        "style": extract_text_style(li),
                    }
                )
        elif tag == "table":
            header_row = [th.get_text(strip=True) for th in child.select("thead th")]
            body_rows = [
                [td.get_text(strip=True) for td in tr.find_all(["td", "th"])]
                for tr in child.select("tbody tr")
            ]
            items.append({"type": "table", "header": header_row, "rows": body_rows})
    return items


def parse_stats(section) -> list[dict]:
    stats_el = section.find("stats", recursive=False)
    if not stats_el:
        return []
    result = []
    for stat_el in stats_el.find_all("stat", recursive=False)[:3]:
        value_el = stat_el.find("statvalue", recursive=False)
        label_el = stat_el.find("statlabel", recursive=False)
        icon_el = stat_el.find("icon", recursive=False)
        value = value_el.get_text(strip=True) if value_el else ""
        label = label_el.get_text(strip=True) if label_el else ""
        icon = validate_icon(icon_el.get_text(strip=True)) if icon_el else ""
        if value or label:
            result.append(
                {
                    "value": value,
                    "label": label,
                    "icon": icon,
                    "container_style": extract_container_style(stat_el),
                    "value_style": extract_text_style(value_el),
                    "label_style": extract_text_style(label_el),
                }
            )
    return result


def parse_boxes(section) -> list[dict]:
    boxes_el = section.find("boxes", recursive=False)
    if not boxes_el:
        return []
    result = []
    for box_el in boxes_el.find_all("box", recursive=False)[:4]:
        icon_el = box_el.find("icon", recursive=False)
        title_el = box_el.find("boxtitle", recursive=False)
        text_el = box_el.find("boxtext", recursive=False)
        icon = validate_icon(icon_el.get_text(strip=True)) if icon_el else ""
        title = title_el.get_text(strip=True) if title_el else ""
        text = text_el.get_text(strip=True) if text_el else ""
        if title:
            result.append(
                {
                    "icon": icon,
                    "title": title,
                    "text": text,
                    "container_style": extract_container_style(box_el),
                    "title_style": extract_text_style(title_el),
                    "text_style": extract_text_style(text_el),
                }
            )
    return result


def parse_columns(section) -> list[dict]:
    columns_el = section.find("columns", recursive=False)
    if not columns_el:
        return []
    return [
        {"items": parse_content_items(col), "container_style": extract_container_style(col)}
        for col in columns_el.find_all("column", recursive=False)[:2]
    ]


def parse_agenda(section) -> list[dict]:
    agenda_el = section.find("agenda", recursive=False)
    if not agenda_el:
        return []
    result = []
    for item_el in agenda_el.find_all("agendaitem", recursive=False)[:6]:
        num_el = item_el.find("agendanum", recursive=False)
        title_el = item_el.find("agendatitle", recursive=False)
        num = num_el.get_text(strip=True) if num_el else ""
        title = title_el.get_text(strip=True) if title_el else ""
        if title:
            result.append(
                {
                    "num": num,
                    "title": title,
                    "container_style": extract_container_style(item_el),
                    "title_style": extract_text_style(title_el),
                }
            )
    return result


def parse_theme(soup) -> str:
    """Liefert den rohen `data-theme`-Wert — kann ein Preset-Key (neue Decks, siehe `PRESET_SPECS`)
    oder ein altes Theme-Key (`THEME_PALETTES`) sein; `resolve_palette` löst beides auf."""
    theme_el = soup.find(attrs={"data-theme": True})
    theme = (theme_el.get("data-theme") or "").strip().lower() if theme_el else ""
    if theme in PRESET_SPECS or theme in THEME_PALETTES:
        return theme
    return DEFAULT_THEME


def parse_slides(html: str) -> list[dict]:
    """`<div data-theme="…"><section class="slide" data-layout="…">…</section>…</div>` → Folien-Liste."""
    soup = BeautifulSoup(f"<div>{html}</div>", "html.parser")
    theme = parse_theme(soup)
    slides = []
    for section in soup.select("section"):
        layout = (section.get("data-layout") or "").strip().lower()
        if layout not in ALLOWED_LAYOUTS:
            layout = "content"

        items = parse_content_items(section)
        slide = {"layout": layout, "theme": theme, "items": items}
        if layout in ("title", "section"):
            slide["section_style"] = extract_container_style(section)

        if layout == "stats":
            stats = parse_stats(section)
            if not stats:
                slide["layout"] = "content"
            else:
                slide["stats"] = stats
                slide["group_style"] = extract_container_style(section.find("stats", recursive=False))
        elif layout == "twocol":
            columns = parse_columns(section)
            if len(columns) < 2:
                slide["layout"] = "content"
            else:
                slide["columns"] = columns
                slide["group_style"] = extract_container_style(section.find("columns", recursive=False))
        elif layout == "agenda":
            agenda_items = parse_agenda(section)
            if not agenda_items:
                slide["layout"] = "content"
            else:
                slide["agenda"] = agenda_items
                slide["group_style"] = extract_container_style(section.find("agenda", recursive=False))
        elif layout == "boxes":
            boxes = parse_boxes(section)
            if len(boxes) < 2:
                slide["layout"] = "content"
            else:
                slide["boxes"] = boxes
                slide["group_style"] = extract_container_style(section.find("boxes", recursive=False))

        if slide["items"] or slide.get("stats") or slide.get("columns") or slide.get("agenda") or slide.get("boxes"):
            slides.append(slide)
    return slides


def add_background(slide, prs, dark: bool, palette: dict) -> None:
    if not dark:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = COLOR_BG_WHITE
        return
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    rect.line.fill.background()
    rect.shadow.inherit = False
    rect.fill.gradient()
    stops = rect.fill.gradient_stops
    stops[0].color.rgb = palette["bg_dark_from"]
    stops[0].position = 0.0
    stops[1].color.rgb = palette["bg_dark_to"]
    stops[1].position = 1.0
    rect.fill.gradient_angle = 135


def add_glass_blob(slide, *, left_in: float, top_in: float, size_in: float, palette: dict) -> None:
    """Glasiger Tropfen mit hellem Rand-Glanzlicht (Tech-Preset) — Annäherung an eine echte radiale
    Glanzlicht-Form über einen 2-Stopp-Linear-Gradient (python-pptx kennt keine radialen Gradients
    ohne XML-Hacking), Winkel passend zum Hintergrund-Gradient (siehe `add_background`)."""
    blob = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left_in), Inches(top_in), Inches(size_in), Inches(size_in))
    blob.line.fill.background()
    blob.shadow.inherit = False
    blob.fill.gradient()
    stops = blob.fill.gradient_stops
    stops[0].color.rgb = palette["accent_on_dark"]
    stops[0].position = 0.0
    stops[1].color.rgb = blend_colors(palette["accent"], palette["bg_dark_to"], 0.55)
    stops[1].position = 1.0
    blob.fill.gradient_angle = 135


def add_accent_spine(slide, prs, palette: dict) -> None:
    """Dünne farbige Kante am linken Folienrand — gibt hellen Layouts einen Marken-Akzent."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.14), prs.slide_height)
    bar.line.fill.background()
    bar.shadow.inherit = False
    bar.fill.solid()
    bar.fill.fore_color.rgb = palette["accent"]


def add_paragraph(
    text_frame, text: str, *, first: bool, size, bold: bool, color, alignment=None, style: dict | None = None, font_name: str | None = None,
) -> None:
    style = style or {}
    paragraph = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
    paragraph.space_after = Pt(10)
    resolved_alignment = style.get("alignment", alignment)
    if resolved_alignment is not None:
        paragraph.alignment = resolved_alignment
    run = paragraph.add_run()
    run.text = text
    factor = style.get("size_factor")
    run.font.size = Pt(size.pt * factor) if factor is not None else size
    run.font.bold = style.get("bold", bold)
    run.font.italic = style.get("italic")
    run.font.underline = style.get("underline")
    run.font.name = font_name or FONT_FAMILY
    run.font.color.rgb = style.get("color", color)


def add_bullet_paragraph(
    text_frame, text: str, *, first: bool, bullet: str, accent: RGBColor, text_color: RGBColor, style: dict | None = None
) -> None:
    """Bullet-Glyphe farbig (Akzent), Text in normaler Textfarbe — Pendant zu CSS `li::marker{color:accent}`."""
    style = style or {}
    paragraph = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
    paragraph.space_after = Pt(10)
    resolved_alignment = style.get("alignment")
    if resolved_alignment is not None:
        paragraph.alignment = resolved_alignment
    bullet_run = paragraph.add_run()
    bullet_run.text = bullet
    bullet_run.font.size = PT_BODY
    bullet_run.font.bold = True
    bullet_run.font.name = FONT_FAMILY
    bullet_run.font.color.rgb = accent
    text_run = paragraph.add_run()
    text_run.text = text
    factor = style.get("size_factor")
    text_run.font.size = Pt(PT_BODY.pt * factor) if factor is not None else PT_BODY
    text_run.font.bold = style.get("bold", False)
    text_run.font.italic = style.get("italic")
    text_run.font.underline = style.get("underline")
    text_run.font.name = FONT_FAMILY
    text_run.font.color.rgb = style.get("color", text_color)


VALIGN_MAP = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}


def _render_title_text_items(text_frame, items: list[dict], palette: dict, *, first: bool) -> bool:
    """Gemeinsamer Render-Loop für `h1`/`h2`/`subtitle`/`p`/`li` der Titelfolie — von
    `render_title_slide` sowohl für den unboxed Fall (alle Items in EINER Box) als auch für die
    Restelemente nach einer geboxten `h1` (eigene Box darunter) genutzt. Textfarben hängen von
    `title_treatment` ab: dunkler Cover-Grund → helle Schrift (Standard), `editorial-light`
    (Soft/Minimal) → heller Grund → dunkle Schrift."""
    light_bg = palette.get("title_treatment") == "editorial-light"
    h1_color = COLOR_TEXT_DARK if light_bg else COLOR_TEXT_LIGHT
    h2_color = palette["accent"] if light_bg else palette["accent_on_dark"]
    subtitle_color = COLOR_STAT_LABEL if light_bg else COLOR_SUBTITLE_LIGHT
    body_color = COLOR_TEXT_DARK if light_bg else COLOR_TEXT_LIGHT
    bullet_accent = palette["accent"] if light_bg else palette["accent_on_dark"]
    heading_font = palette.get("heading_font", FONT_FAMILY)
    for item in items:
        kind = item["type"]
        if not item.get("text"):
            continue
        if kind == "h1":
            add_paragraph(
                text_frame, item["text"], first=first, size=PT_H1, bold=False, color=h1_color, style=item.get("style"), font_name=heading_font,
            )
        elif kind == "h2":
            add_paragraph(
                text_frame, item["text"], first=first, size=PT_H2, bold=True, color=h2_color, style=item.get("style"), font_name=heading_font,
            )
        elif kind == "subtitle":
            add_paragraph(
                text_frame, item["text"], first=first, size=PT_SUBTITLE, bold=False, color=subtitle_color, style=item.get("style")
            )
        elif kind == "p":
            add_paragraph(text_frame, item["text"], first=first, size=PT_BODY, bold=False, color=body_color, style=item.get("style"))
        elif kind == "li":
            bullet = f"{item['index']}. " if item.get("ordered") else "•  "
            add_bullet_paragraph(
                text_frame,
                item["text"],
                first=first,
                bullet=bullet,
                accent=bullet_accent,
                text_color=body_color,
                style=item.get("style"),
            )
        else:
            continue
        first = False
    return first


def render_title_slide(slide, items: list[dict], palette: dict, section_style: dict | None = None) -> None:
    """`title` — die einzige verbleibende Cover-Folie: Gradient + dezentes Deko-Rondell, Text
    linksbündig in einer schmaleren Box statt voller Breite (Editorial-/Poster-Look).
    Trägt `<h1>` ein `data-color` (siehe `extract_text_style`), wird der Titel stattdessen als
    eigene farbige Box dargestellt — feste Höhe, da python-pptx die tatsächliche Texthöhe nicht
    kennt (gleiche Pragmatik wie bei `callout`, siehe `render_content_slide`)."""
    section_style = section_style or {}
    # Dezentes Deko-Element nur beim dunklen Cover-Look — `editorial-light` (Soft/Minimal) bleibt
    # bewusst ohne zusätzliches Deko-Element (ruhigerer, "editorial" Eindruck).
    if palette.get("title_treatment") != "editorial-light":
        if palette.get("decoration_style") == "blob":
            # Tech-Preset: organische, glasige Tropfenform mit hellem Rand-Glanzlicht statt eines
            # schlichten Kreises — zwei Tropfen (gross unten rechts, klein separat darüber), an ein
            # Referenzbild angelehnt.
            add_glass_blob(slide, left_in=SLIDE_WIDTH_IN - 3.4, top_in=SLIDE_HEIGHT_IN - 3.6, size_in=4.8, palette=palette)
            add_glass_blob(slide, left_in=SLIDE_WIDTH_IN - 4.6, top_in=SLIDE_HEIGHT_IN - 5.4, size_in=1.3, palette=palette)
        else:
            decoration_color = blend_colors(palette["accent_on_dark"], palette["bg_dark_to"], 0.3)
            deco_size_in = 4.4
            decoration = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(SLIDE_WIDTH_IN - 3.1),
                Inches(SLIDE_HEIGHT_IN - 3.1),
                Inches(deco_size_in),
                Inches(deco_size_in),
            )
            decoration.line.fill.background()
            decoration.shadow.inherit = False
            decoration.fill.solid()
            decoration.fill.fore_color.rgb = decoration_color

    box_width = CONTENT_WIDTH_IN * 0.62
    heading_item = next((i for i in items if i["type"] == "h1"), None)
    heading_color_key = (heading_item or {}).get("style", {}).get("color_key") if heading_item else None

    if heading_item and heading_color_key:
        heading_palette = THEME_PALETTES.get(heading_color_key, palette)
        radius = RADIUS_ADJUSTMENTS.get(heading_item["style"].get("radius") or "", corner(palette, 0.12))
        title_box_height = 1.6
        title_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(MARGIN_X_IN), Inches(MARGIN_Y_IN), Inches(box_width), Inches(title_box_height)
        )
        title_box.line.fill.background()
        title_box.shadow.inherit = False
        title_box.fill.solid()
        title_box.fill.fore_color.rgb = heading_palette["accent"]
        title_box.adjustments[0] = radius
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        title_tf.margin_left = Inches(0.4)
        title_tf.margin_right = Inches(0.4)
        _render_title_text_items(title_tf, [heading_item], palette, first=True)

        other_items = [i for i in items if i is not heading_item]
        remaining_top = MARGIN_Y_IN + title_box_height + 0.3
        remaining_height = SLIDE_HEIGHT_IN - MARGIN_Y_IN - remaining_top
        if other_items and remaining_height > 0.3:
            box = slide.shapes.add_textbox(
                Inches(MARGIN_X_IN), Inches(remaining_top), Inches(box_width), Inches(remaining_height)
            )
            tf = box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.TOP
            _render_title_text_items(tf, other_items, palette, first=True)
        return

    box = slide.shapes.add_textbox(
        Inches(MARGIN_X_IN), Inches(MARGIN_Y_IN), Inches(box_width), Inches(CONTENT_HEIGHT_IN)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = VALIGN_MAP.get(section_style.get("valign"), MSO_ANCHOR.MIDDLE)
    _render_title_text_items(tf, items, palette, first=True)


def render_section_slide(slide, items: list[dict], palette: dict, section_style: dict | None = None) -> None:
    """`section` — heller Hintergrund wie alle Nicht-Cover-Layouts, Titel als flache Akzent-Box (kein Gradient mehr)."""
    section_style = section_style or {}
    heading = next((i for i in items if i["type"] == "h1" and i.get("text")), None)
    if not heading:
        return
    box_width = CONTENT_WIDTH_IN * 0.7
    box_height = 2.2
    valign = section_style.get("valign")
    if valign == "top":
        y = MARGIN_Y_IN
    elif valign == "bottom":
        y = SLIDE_HEIGHT_IN - MARGIN_Y_IN - box_height
    else:
        y = (SLIDE_HEIGHT_IN - box_height) / 2
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches((SLIDE_WIDTH_IN - box_width) / 2),
        Inches(y),
        Inches(box_width),
        Inches(box_height),
    )
    box.line.fill.background()
    box.shadow.inherit = False
    box.fill.solid()
    box.fill.fore_color.rgb = palette["accent"]
    box.adjustments[0] = corner(palette, 0.12)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.4)
    tf.margin_right = Inches(0.4)
    add_paragraph(
        tf, heading["text"], first=True, size=PT_H1, bold=True, color=COLOR_TEXT_LIGHT, alignment=PP_ALIGN.CENTER, style=heading.get("style"),
        font_name=palette.get("heading_font", FONT_FAMILY),
    )


def render_content_slide(slide, items: list[dict], palette: dict, dark: bool = False) -> None:
    """`content` — Titel oben/fix, Inhalt darunter oben ausgerichtet (nicht vertikal zentriert)."""
    heading = next((i for i in items if i["type"] in ("h1", "h2")), None)
    callout = next((i for i in items if i["type"] == "callout" and i.get("text")), None)
    body_items = [i for i in items if i["type"] not in ("h1", "h2", "callout")]
    heading_color = palette["accent_on_dark"] if dark else palette["accent"]
    body_color = COLOR_TEXT_LIGHT if dark else COLOR_TEXT_DARK
    bullet_accent = palette["accent_on_dark"] if dark else palette["accent"]

    top = MARGIN_Y_IN
    if heading:
        heading_box = slide.shapes.add_textbox(
            Inches(MARGIN_X_IN), Inches(top), Inches(CONTENT_WIDTH_IN), Inches(0.9)
        )
        tf = heading_box.text_frame
        tf.word_wrap = True
        size = PT_H1 if heading["type"] == "h1" else PT_H2
        add_paragraph(
            tf, heading["text"], first=True, size=size, bold=True, color=heading_color, style=heading.get("style"),
            font_name=palette.get("heading_font", FONT_FAMILY),
        )
        top += 1.1

    callout_height = 1.3 if callout else 0.0
    body_bottom_gap = (callout_height + 0.3) if callout else 0.0
    body_height = max(SLIDE_HEIGHT_IN - MARGIN_Y_IN - top - body_bottom_gap, 0.8)

    if body_items:
        box = slide.shapes.add_textbox(Inches(MARGIN_X_IN), Inches(top), Inches(CONTENT_WIDTH_IN), Inches(body_height))
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        first = True
        for item in body_items:
            kind = item["type"]
            if not item.get("text"):
                continue
            if kind == "subtitle":
                add_paragraph(tf, item["text"], first=first, size=PT_SUBTITLE, bold=False, color=body_color, style=item.get("style"))
            elif kind == "p":
                add_paragraph(tf, item["text"], first=first, size=PT_BODY, bold=False, color=body_color, style=item.get("style"))
            elif kind == "li":
                bullet = f"{item['index']}. " if item.get("ordered") else "•  "
                add_bullet_paragraph(
                    tf,
                    item["text"],
                    first=first,
                    bullet=bullet,
                    accent=bullet_accent,
                    text_color=body_color,
                    style=item.get("style"),
                )
            else:
                continue
            first = False

    if callout:
        callout_style = callout.get("container_style", {})
        callout_top = SLIDE_HEIGHT_IN - MARGIN_Y_IN - callout_height
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(MARGIN_X_IN), Inches(callout_top), Inches(CONTENT_WIDTH_IN), Inches(callout_height)
        )
        box.shadow.inherit = False
        box.line.fill.background()
        box.fill.solid()
        box.fill.fore_color.rgb = THEME_PALETTES[callout_style["color_key"]]["accent"] if "color_key" in callout_style else palette["accent"]
        box.adjustments[0] = callout_style.get("radius", corner(palette, 0.08))
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.3)
        tf.margin_right = Inches(0.3)
        # Kein Emoji-Icon im dunklen Tech-Deck (emoji-frei); helle Presets behalten ihr Callout-Icon.
        callout_text = f"{callout['icon']}  {callout['text']}" if callout.get("icon") and not dark else callout["text"]
        add_paragraph(
            tf, callout_text, first=True, size=PT_BODY, bold=True, color=COLOR_TEXT_LIGHT, alignment=callout_style.get("alignment")
        )


def render_agenda_slide(slide, items: list[dict], agenda_items: list[dict], palette: dict, group_style: dict, dark: bool = False) -> None:
    heading = next((i for i in items if i["type"] == "h2"), None)
    heading_color = palette["accent_on_dark"] if dark else palette["accent"]
    title_color = COLOR_TEXT_LIGHT if dark else COLOR_TEXT_DARK
    top = MARGIN_Y_IN
    if heading:
        heading_box = slide.shapes.add_textbox(
            Inches(MARGIN_X_IN), Inches(top), Inches(CONTENT_WIDTH_IN), Inches(0.9)
        )
        tf = heading_box.text_frame
        tf.word_wrap = True
        add_paragraph(
            tf, heading["text"], first=True, size=PT_H2, bold=True, color=heading_color, style=heading.get("style"),
            font_name=palette.get("heading_font", FONT_FAMILY),
        )
        top += 1.1

    badge_size = 0.7
    pad = 0.18
    gap = 0.2
    row_height = badge_size + 2 * pad
    available_height = max(SLIDE_HEIGHT_IN - MARGIN_Y_IN - top, row_height)
    max_rows = max(int((available_height + gap) // (row_height + gap)), 1)
    used_rows = min(len(agenda_items), max_rows)
    list_height = used_rows * row_height + max(used_rows - 1, 0) * gap
    valign = group_style.get("valign", "top")
    if valign == "middle":
        top += max((available_height - list_height) / 2, 0)
    elif valign == "bottom":
        top += max(available_height - list_height, 0)

    for i, entry in enumerate(agenda_items[:max_rows]):
        item_style = entry.get("container_style", {})
        row_top = top + i * (row_height + gap)
        badge_top = row_top + pad
        badge_left = MARGIN_X_IN + pad

        panel = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(MARGIN_X_IN), Inches(row_top), Inches(CONTENT_WIDTH_IN), Inches(row_height)
        )
        panel.line.fill.background()
        panel.shadow.inherit = False
        panel.fill.solid()
        panel_accent = THEME_PALETTES[item_style["color_key"]]["accent"] if "color_key" in item_style else palette["accent"]
        panel.fill.fore_color.rgb = dark_card_fill(palette, 0.1) if dark else mix_with_white(panel_accent, 0.05)
        panel.adjustments[0] = item_style.get("radius", corner(palette, 0.12))

        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(badge_left), Inches(badge_top), Inches(badge_size), Inches(badge_size)
        )
        badge.line.fill.background()
        badge.shadow.inherit = False
        badge.fill.solid()
        badge.fill.fore_color.rgb = panel_accent
        badge.adjustments[0] = corner(palette, 0.2)
        tf = badge.text_frame
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        add_paragraph(tf, entry["num"], first=True, size=Pt(20), bold=True, color=COLOR_TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

        title_left = badge_left + badge_size + 0.3
        title_box = slide.shapes.add_textbox(
            Inches(title_left), Inches(badge_top), Inches(MARGIN_X_IN + CONTENT_WIDTH_IN - pad - title_left), Inches(badge_size)
        )
        tf2 = title_box.text_frame
        tf2.word_wrap = True
        tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
        add_paragraph(
            tf2,
            entry["title"],
            first=True,
            size=Pt(20),
            bold=False,
            color=title_color,
            alignment=item_style.get("alignment"),
            style=entry.get("title_style"),
        )


def render_table_slide(slide, items: list[dict], palette: dict, dark: bool = False) -> None:
    heading = next((i for i in items if i["type"] == "h2"), None)
    table_item = next((i for i in items if i["type"] == "table"), None)
    heading_color = palette["accent_on_dark"] if dark else palette["accent"]

    heading_box = slide.shapes.add_textbox(
        Inches(MARGIN_X_IN), Inches(MARGIN_Y_IN), Inches(CONTENT_WIDTH_IN), Inches(0.8)
    )
    tf = heading_box.text_frame
    tf.word_wrap = True
    if heading:
        add_paragraph(
            tf, heading["text"], first=True, size=PT_H2, bold=True, color=heading_color, font_name=palette.get("heading_font", FONT_FAMILY),
        )

    if not table_item:
        return
    header = table_item.get("header") or []
    rows = table_item.get("rows") or []
    col_count = max([len(header)] + [len(r) for r in rows], default=1) or 1
    row_count = len(rows) + (1 if header else 0)
    if row_count == 0:
        return

    table_top = MARGIN_Y_IN + 1.0
    table_height = max(SLIDE_HEIGHT_IN - MARGIN_Y_IN - table_top, 0.5)
    graphic_frame = slide.shapes.add_table(
        row_count, col_count, Inches(MARGIN_X_IN), Inches(table_top), Inches(CONTENT_WIDTH_IN), Inches(table_height)
    )
    table = graphic_frame.table
    if dark:
        header_fill = dark_card_fill(palette, 0.2)
        body_fill = palette["bg_dark_to"]
        band_fill = dark_card_fill(palette, 0.07)
        header_text = palette["accent_on_dark"]
        body_text = COLOR_TEXT_LIGHT
    else:
        header_fill = mix_with_white(palette["accent"], 0.16)
        body_fill = COLOR_BG_WHITE
        band_fill = mix_with_white(palette["accent"], 0.06)
        header_text = COLOR_TEXT_DARK
        body_text = COLOR_TEXT_DARK

    def style_cell(cell, text: str, *, header_row: bool, banded: bool) -> None:
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill if header_row else (band_fill if banded else body_fill)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = PT_TABLE
                run.font.bold = header_row
                run.font.name = FONT_FAMILY
                run.font.color.rgb = header_text if header_row else body_text

    row_offset = 0
    if header:
        for col, text in enumerate(header):
            if col < col_count:
                style_cell(table.cell(0, col), text, header_row=True, banded=False)
        row_offset = 1
    for r, row in enumerate(rows):
        for col, text in enumerate(row):
            if col < col_count:
                style_cell(table.cell(r + row_offset, col), text, header_row=False, banded=r % 2 == 1)


def render_stats_slide(slide, items: list[dict], stats: list[dict], palette: dict, group_style: dict, dark: bool = False) -> None:
    heading = next((i for i in items if i["type"] == "h2"), None)
    heading_color = palette["accent_on_dark"] if dark else palette["accent"]
    top = MARGIN_Y_IN
    if heading:
        heading_box = slide.shapes.add_textbox(
            Inches(MARGIN_X_IN), Inches(top), Inches(CONTENT_WIDTH_IN), Inches(0.9)
        )
        tf = heading_box.text_frame
        tf.word_wrap = True
        add_paragraph(
            tf, heading["text"], first=True, size=PT_H2, bold=True, color=heading_color, style=heading.get("style"),
            font_name=palette.get("heading_font", FONT_FAMILY),
        )
        top += 1.1

    count = max(len(stats), 1)
    gap = 0.4
    box_width = (CONTENT_WIDTH_IN - gap * (count - 1)) / count
    available_height = max(SLIDE_HEIGHT_IN - MARGIN_Y_IN - top, 1.0)
    box_height = min(available_height, 3.2)
    valign = group_style.get("valign", "middle")
    if valign == "top":
        box_top = top
    elif valign == "bottom":
        box_top = top + max(available_height - box_height, 0)
    else:
        box_top = top + max((available_height - box_height) / 2, 0)

    for i, stat in enumerate(stats):
        item_style = stat.get("container_style", {})
        override = THEME_PALETTES.get(item_style["color_key"]) if "color_key" in item_style else None
        accent = override["accent"] if override else palette["accent"]
        left = MARGIN_X_IN + i * (box_width + gap)
        radius = item_style.get("radius", corner(palette, 0.08))
        alignment = item_style.get("alignment", PP_ALIGN.CENTER)

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(box_top), Inches(box_width), Inches(box_height)
        )
        box.line.fill.background()
        box.shadow.inherit = False
        box.fill.solid()
        box.adjustments[0] = radius

        if dark:
            # Dunkle Karte + Teal-Hairline oben + Glow-Ring hinter dem Wert. Text liegt in einer
            # SEPARATEN Box ÜBER Karte und Ring (in python-pptx liegen später erzeugte Shapes oben).
            value_color = override["accent_on_dark"] if override else palette["accent_on_dark"]
            box.fill.fore_color.rgb = dark_card_fill(palette, 0.14)
            add_card_hairline(slide, left_in=left, top_in=box_top, width_in=box_width, color=accent)
            add_glow_ring(
                slide, palette, cx_in=left + box_width / 2, cy_in=box_top + box_height * 0.4,
                radius_in=min(box_width, box_height) * 0.26,
            )
            tbox = slide.shapes.add_textbox(
                Inches(left + 0.15), Inches(box_top), Inches(box_width - 0.3), Inches(box_height)
            )
            tf = tbox.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            # Kein Emoji-Icon im dunklen Tech-Deck (Präsentationen bleiben emoji-frei).
            add_paragraph(
                tf, stat["value"], first=True, size=PT_STAT_VALUE, bold=True, color=value_color, alignment=alignment, style=stat.get("value_style")
            )
            add_paragraph(
                tf, stat["label"], first=False, size=PT_STAT_LABEL, bold=False, color=COLOR_MUTED_ON_DARK,
                alignment=alignment, style=stat.get("label_style"),
            )
            continue

        box.fill.fore_color.rgb = mix_with_white(accent, 0.12)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        first = True
        if stat.get("icon"):
            add_paragraph(tf, stat["icon"], first=first, size=Pt(32), bold=False, color=accent, alignment=alignment)
            first = False
        add_paragraph(
            tf, stat["value"], first=first, size=PT_STAT_VALUE, bold=True, color=accent, alignment=alignment, style=stat.get("value_style")
        )
        add_paragraph(
            tf,
            stat["label"],
            first=False,
            size=PT_STAT_LABEL,
            bold=False,
            color=COLOR_STAT_LABEL,
            alignment=alignment,
            style=stat.get("label_style"),
        )


def render_twocol_slide(slide, items: list[dict], columns: list[dict], palette: dict, dark: bool = False) -> None:
    heading = next((i for i in items if i["type"] == "h2"), None)
    heading_color = palette["accent_on_dark"] if dark else palette["accent"]
    body_color = COLOR_TEXT_LIGHT if dark else COLOR_TEXT_DARK
    top = MARGIN_Y_IN
    if heading:
        heading_box = slide.shapes.add_textbox(
            Inches(MARGIN_X_IN), Inches(top), Inches(CONTENT_WIDTH_IN), Inches(0.9)
        )
        tf = heading_box.text_frame
        tf.word_wrap = True
        add_paragraph(
            tf, heading["text"], first=True, size=PT_H2, bold=True, color=heading_color, style=heading.get("style"),
            font_name=palette.get("heading_font", FONT_FAMILY),
        )
        top += 1.1

    gap = 0.6
    col_width = (CONTENT_WIDTH_IN - gap) / 2
    col_height = max(SLIDE_HEIGHT_IN - MARGIN_Y_IN - top, 1.0)

    for i, column in enumerate(columns[:2]):
        col_style = column.get("container_style", {})
        override = THEME_PALETTES.get(col_style["color_key"]) if "color_key" in col_style else None
        accent = override["accent"] if override else palette["accent"]
        # Spalten-Akzent (h2/Bullet) auf Dunkel: helle Teal-Variante, damit er auf der dunklen Karte trägt.
        col_accent = (override["accent_on_dark"] if override else palette["accent_on_dark"]) if dark else accent
        left = MARGIN_X_IN + i * (col_width + gap)

        panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(col_width), Inches(col_height))
        panel.line.fill.background()
        panel.shadow.inherit = False
        panel.fill.solid()
        panel.fill.fore_color.rgb = dark_card_fill(palette, 0.1) if dark else mix_with_white(accent, 0.06)
        panel.adjustments[0] = col_style.get("radius", corner(palette, 0.06))

        pad = 0.35
        box = slide.shapes.add_textbox(
            Inches(left + pad), Inches(top + pad), Inches(col_width - 2 * pad), Inches(col_height - 2 * pad)
        )
        tf = box.text_frame
        tf.word_wrap = True
        first = True
        for item in column.get("items", []):
            kind = item["type"]
            if not item.get("text"):
                continue
            if kind == "h2":
                add_paragraph(
                    tf, item["text"], first=first, size=PT_COL_H2, bold=True, color=col_accent,
                    alignment=col_style.get("alignment"), style=item.get("style"),
                )
            elif kind == "p":
                add_paragraph(
                    tf, item["text"], first=first, size=PT_BODY, bold=False, color=body_color,
                    alignment=col_style.get("alignment"), style=item.get("style"),
                )
            elif kind == "li":
                bullet = f"{item['index']}. " if item.get("ordered") else "•  "
                add_bullet_paragraph(
                    tf,
                    item["text"],
                    first=first,
                    bullet=bullet,
                    accent=col_accent,
                    text_color=body_color,
                    style=item.get("style"),
                )
            else:
                continue
            first = False


def render_boxes_slide(slide, items: list[dict], boxes: list[dict], palette: dict, dark: bool = False) -> None:
    """`boxes` — Vorteile/Schritte/Features als Karten. Helles Deck: flache Vollfarb-Karten mit
    zyklisch wechselnden Farben. Dunkles Deck (Tech): dunkel-transluzente Karten mit Teal-Hairline
    oben, Teal-Icon und hellem Text (an das Referenzbild «Everything is information» angelehnt)."""
    heading = next((i for i in items if i["type"] == "h2"), None)
    heading_color = palette["accent_on_dark"] if dark else palette["accent"]
    top = MARGIN_Y_IN
    if heading:
        heading_box = slide.shapes.add_textbox(
            Inches(MARGIN_X_IN), Inches(top), Inches(CONTENT_WIDTH_IN), Inches(0.9)
        )
        tf = heading_box.text_frame
        tf.word_wrap = True
        add_paragraph(
            tf, heading["text"], first=True, size=PT_H2, bold=True, color=heading_color, style=heading.get("style"),
            font_name=palette.get("heading_font", FONT_FAMILY),
        )
        top += 1.1

    count = max(len(boxes), 1)
    gap = 0.4
    box_width = (CONTENT_WIDTH_IN - gap * (count - 1)) / count
    box_height = max(SLIDE_HEIGHT_IN - MARGIN_Y_IN - top, 1.0)
    box_colors = palette["box_colors"]

    for i, item in enumerate(boxes):
        item_style = item.get("container_style", {})
        override = THEME_PALETTES.get(item_style["color_key"]) if "color_key" in item_style else None
        left = MARGIN_X_IN + i * (box_width + gap)
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(box_width), Inches(box_height)
        )
        box.line.fill.background()
        box.shadow.inherit = False
        box.fill.solid()
        box.adjustments[0] = item_style.get("radius", corner(palette, 0.08))

        if dark:
            accent = override["accent"] if override else palette["accent"]
            icon_color = override["accent_on_dark"] if override else palette["accent_on_dark"]
            box.fill.fore_color.rgb = dark_card_fill(palette, 0.1)
            add_card_hairline(slide, left_in=left, top_in=top, width_in=box_width, color=accent)
            title_color = COLOR_TEXT_LIGHT
            text_color = COLOR_MUTED_ON_DARK
        else:
            box.fill.fore_color.rgb = override["accent"] if override else box_colors[i % len(box_colors)]
            icon_color = COLOR_TEXT_LIGHT
            title_color = COLOR_TEXT_LIGHT
            text_color = COLOR_TEXT_LIGHT

        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = Inches(0.25)
        tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.3)
        alignment = item_style.get("alignment")
        first = True
        # Kein Emoji-Icon im dunklen Tech-Deck (emoji-frei); helle Presets behalten ihr Icon.
        if item.get("icon") and not dark:
            add_paragraph(tf, item["icon"], first=first, size=PT_BOX_ICON, bold=False, color=icon_color, alignment=alignment)
            first = False
        add_paragraph(
            tf, item["title"], first=first, size=PT_BOX_TITLE, bold=True, color=title_color, alignment=alignment, style=item.get("title_style")
        )
        if item.get("text"):
            add_paragraph(
                tf, item["text"], first=False, size=PT_BOX_TEXT, bold=False, color=text_color, alignment=alignment, style=item.get("text_style")
            )


def build_presentation(slides: list[dict]) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    blank_layout = prs.slide_layouts[6]

    for index, slide_data in enumerate(slides):
        slide = prs.slides.add_slide(blank_layout)
        layout = slide_data["layout"]
        palette = resolve_palette(slide_data.get("theme", DEFAULT_THEME))
        # Nur die Titelfolie behält den grossflächigen Gradient-Cover-Look (`gradient-cover` ODER
        # `editorial-light` — beide nutzen einen 2-Stop-Verlauf, nur mit anderen Farben, siehe
        # `add_background`) — alle anderen Layouts inkl. `section` sind hell und bekommen den
        # linken Akzent-Rand, sofern das Preset ihn nicht bewusst abschaltet (`accent_spine`).
        is_title = layout == "title"
        dark = is_dark_deck(palette)
        # Dunkles Tech-Deck: jede Folie dunkel (nicht nur Titel). Helle Presets: nur Titel dunkel.
        add_background(slide, prs, is_title or dark, palette)
        # Heller Akzent-Streifen nur auf hellen Decks; das dunkle Deck trägt seine Akzente über
        # Deko-Muster und Teal-Karten, ein Vollflächen-Streifen würde dort fehl wirken.
        if not is_title and not dark and palette.get("accent_spine", True):
            add_accent_spine(slide, prs, palette)
        # Ambiente Deko VOR dem Inhalt einfügen (spätere Shapes liegen in python-pptx oben), damit
        # Text/Karten darüber liegen. EIN Dot-Band pro Inhalts-Folie, Position rotiert über den
        # Folien-Index (`index` = gemeinsamer Schlüssel mit der Preview-CSS, siehe `buildPptxSlideSrcDoc`).
        if dark and not is_title:
            placement = DOT_BAND_PLACEMENTS[index % len(DOT_BAND_PLACEMENTS)]
            add_dot_band(slide, palette, placement)

        if layout == "table":
            render_table_slide(slide, slide_data["items"], palette, dark)
        elif layout == "stats":
            render_stats_slide(slide, slide_data["items"], slide_data.get("stats", []), palette, slide_data.get("group_style", {}), dark)
        elif layout == "twocol":
            render_twocol_slide(slide, slide_data["items"], slide_data.get("columns", []), palette, dark)
        elif layout == "agenda":
            render_agenda_slide(slide, slide_data["items"], slide_data.get("agenda", []), palette, slide_data.get("group_style", {}), dark)
        elif layout == "boxes":
            render_boxes_slide(slide, slide_data["items"], slide_data.get("boxes", []), palette, dark)
        elif layout == "content":
            render_content_slide(slide, slide_data["items"], palette, dark)
        elif layout == "section":
            render_section_slide(slide, slide_data["items"], palette, slide_data.get("section_style", {}))
        else:
            render_title_slide(slide, slide_data["items"], palette, slide_data.get("section_style", {}))

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


# ---------------------------------------------------------------------------
# Word (.docx) — feste Formatierung (kein Theme). Gespiegelte Konstanten zu
# `src/features/chat/constants/wordDocStyle.ts` (Single Source of Truth dort). Bei einer
# Änderung dort müssen diese Werte mitgezogen werden.
# ---------------------------------------------------------------------------
WORD_FONT_NAME = "Arial"
WORD_HEADING_RGB = DocxRGB(0x1F, 0x4E, 0x79)  # Dunkelblau für H1 UND H2 (identisch)
WORD_BODY_RGB = DocxRGB(0x00, 0x00, 0x00)
WORD_HEADER_FILL = "F1F4F8"  # helles Tabellenkopf-Grau (entspricht der CSS-Vorschau)
WORD_HEADER_TEXT_RGB = DocxRGB(0x80, 0x80, 0x80)  # Kopfzeilen-Titel in Grau
WORD_HEADER_LINE_HEX = "C8C8C8"  # Kopfzeilen-Unterlinie
WORD_COVER_SUBTITLE_RGB = DocxRGB(0x44, 0x44, 0x44)  # Untertitel auf dem Titelblatt
WORD_COVER_META_RGB = DocxRGB(0x55, 0x55, 0x55)  # Autor + Datum unten auf dem Titelblatt

WORD_SPEC = {
    "body": {"size": 11.0, "line": 1.15, "after": 8.0},
    "h1": {"size": 16.0, "before": 16.0, "after": 4.0},
    "h2": {"size": 13.0, "before": 10.0, "after": 2.0},
    "h3": {"size": 11.5, "before": 8.0, "after": 2.0},
    "list": {"item_gap": 2.0, "after": 8.0},
    "table": {"size": 11.0},
    "cover": {"title_size": 28.0, "subtitle_size": 15.0, "meta_size": 11.0, "top_space": 96.0},
    "header": {"size": 9.0},
}
MAX_WORD_BLOCKS = 5000


def _set_run_font(run, *, bold: bool, color: DocxRGB, size_pt: float) -> None:
    run.font.name = WORD_FONT_NAME
    # Schriftart auch für Komplex-/Ostasiatische Zeichen erzwingen (sonst greift Calibri-Fallback).
    rpr = run._element.get_or_add_rPr()
    rfonts = rpr.find(qn("w:rFonts"))
    if rfonts is None:
        rfonts = OxmlElement("w:rFonts")
        rpr.append(rfonts)
    for attr in ("w:ascii", "w:hAnsi", "w:cs"):
        rfonts.set(qn(attr), WORD_FONT_NAME)
    run.font.bold = bold
    run.font.size = DocxPt(size_pt)
    run.font.color.rgb = color


def _add_runs(paragraph, text: str, *, bold: bool, color: DocxRGB, size_pt: float) -> None:
    """`**fett**` → fette Runs; `\\n` → echter Zeilenumbruch. Gleiche Inline-Regel wie `wordInlineToHtml`."""
    lines = text.split("\n")
    for line_index, line in enumerate(lines):
        if line_index > 0:
            paragraph.add_run().add_break()
        # an **…** aufteilen; ungerade Segmente sind fett
        for seg_index, segment in enumerate(line.split("**")):
            if segment == "":
                continue
            run = paragraph.add_run(segment)
            _set_run_font(run, bold=bold or (seg_index % 2 == 1), color=color, size_pt=size_pt)


def _style_paragraph(paragraph, *, before_pt: float, after_pt: float, line: float) -> None:
    pf = paragraph.paragraph_format
    pf.space_before = DocxPt(before_pt)
    pf.space_after = DocxPt(after_pt)
    pf.line_spacing = line


def _set_cell_background(cell, hex_fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), hex_fill)
    tc_pr.append(shd)


def _heading_spec(level: int) -> tuple[dict, DocxRGB]:
    if level <= 1:
        return WORD_SPEC["h1"], WORD_HEADING_RGB
    if level == 2:
        return WORD_SPEC["h2"], WORD_HEADING_RGB
    return WORD_SPEC["h3"], WORD_BODY_RGB  # Ebene ≥3: fett, aber in Textfarbe


def _add_heading_block(doc, block: dict) -> None:
    try:
        level = int(block.get("level", 1))
    except (TypeError, ValueError):
        level = 1
    spec, color = _heading_spec(level)
    p = doc.add_paragraph()
    _style_paragraph(p, before_pt=spec["before"], after_pt=spec["after"], line=1.2)
    _add_runs(p, str(block.get("text", "")), bold=True, color=color, size_pt=spec["size"])


def _add_paragraph_block(doc, block: dict) -> None:
    body = WORD_SPEC["body"]
    p = doc.add_paragraph()
    _style_paragraph(p, before_pt=0.0, after_pt=body["after"], line=body["line"])
    _add_runs(p, str(block.get("text", "")), bold=False, color=WORD_BODY_RGB, size_pt=body["size"])


def _add_list_block(doc, block: dict) -> None:
    items = [str(it) for it in (block.get("items") or []) if str(it).strip()]
    if not items:
        return
    ordered = block.get("ordered") is True
    style = "List Number" if ordered else "List Bullet"
    body = WORD_SPEC["body"]
    spec = WORD_SPEC["list"]
    for idx, item in enumerate(items):
        try:
            p = doc.add_paragraph(style=style)
        except KeyError:
            # Fallback, falls die Vorlage den Listenstil nicht kennt: manuelles Aufzählungszeichen.
            p = doc.add_paragraph()
            item = f"{idx + 1}. {item}" if ordered else f"•  {item}"
        last = idx == len(items) - 1
        _style_paragraph(p, before_pt=0.0, after_pt=(spec["after"] if last else spec["item_gap"]), line=body["line"])
        _add_runs(p, item, bold=False, color=WORD_BODY_RGB, size_pt=body["size"])


def _add_table_block(doc, block: dict) -> None:
    rows = block.get("rows") or []
    rows = [r for r in rows if isinstance(r, list)]
    if not rows:
        return
    cols = max((len(r) for r in rows), default=0)
    if cols == 0:
        return
    has_header = block.get("header") is True
    size = WORD_SPEC["table"]["size"]
    table = doc.add_table(rows=len(rows), cols=cols)
    table.style = "Table Grid"
    for i, row in enumerate(rows):
        is_header = has_header and i == 0
        for j in range(cols):
            cell = table.cell(i, j)
            value = str(row[j]) if j < len(row) else ""
            para = cell.paragraphs[0]
            _style_paragraph(para, before_pt=0.0, after_pt=0.0, line=1.1)
            _add_runs(para, value, bold=is_header, color=WORD_BODY_RGB, size_pt=size)
            if is_header:
                _set_cell_background(cell, WORD_HEADER_FILL)
    # kleiner Abstand nach der Tabelle (Word-Tabellen tragen keinen eigenen)
    spacer = doc.add_paragraph()
    _style_paragraph(spacer, before_pt=0.0, after_pt=0.0, line=1.0)
    spacer.add_run("").font.size = DocxPt(4)


def _apply_a4_section(section) -> None:
    """A4-Hochformat mit 1-Zoll-Rand (= 72 pt, wie `WORD_DOC_SPEC.marginPt`)."""
    section.page_width = Cm(21.0)
    section.page_height = Cm(29.7)
    section.top_margin = Cm(2.54)
    section.bottom_margin = Cm(2.54)
    section.left_margin = Cm(2.54)
    section.right_margin = Cm(2.54)


def _add_bottom_border(paragraph, *, color_hex: str, size: int) -> None:
    """Untere Linie unter einem Absatz (für Kopfzeile + Titelblatt-Akzentlinie)."""
    p_pr = paragraph._p.get_or_add_pPr()
    p_bdr = OxmlElement("w:pBdr")
    bottom = OxmlElement("w:bottom")
    bottom.set(qn("w:val"), "single")
    bottom.set(qn("w:sz"), str(size))
    bottom.set(qn("w:space"), "4")
    bottom.set(qn("w:color"), color_hex)
    p_bdr.append(bottom)
    p_pr.append(p_bdr)


def _add_cover_page(doc, title: str, subtitle: str, author: str, date: str) -> None:
    """Titelblatt: linksbündiger Titel (dunkelblau) + Untertitel oben; Autor + Datum unten in der
    Fusszeile der Cover-Sektion. Keine Trennlinie."""
    spec = WORD_SPEC["cover"]

    p_title = doc.add_paragraph()
    p_title.alignment = WD_ALIGN_PARAGRAPH.LEFT
    _style_paragraph(p_title, before_pt=spec["top_space"], after_pt=0.0, line=1.2)
    _add_runs(p_title, title, bold=True, color=WORD_HEADING_RGB, size_pt=spec["title_size"])

    if subtitle:
        p_sub = doc.add_paragraph()
        p_sub.alignment = WD_ALIGN_PARAGRAPH.LEFT
        _style_paragraph(p_sub, before_pt=6.0, after_pt=0.0, line=1.3)
        _add_runs(p_sub, subtitle, bold=False, color=WORD_COVER_SUBTITLE_RGB, size_pt=spec["subtitle_size"])

    # Autor + Datum unten: Fusszeile der Cover-Sektion (sitzt im unteren Seitenrand).
    meta_lines = [(author, True), (date, False)]
    meta_lines = [(txt, bold) for (txt, bold) in meta_lines if txt]
    footer = doc.sections[0].footer
    footer.is_linked_to_previous = False
    for idx, (txt, bold) in enumerate(meta_lines):
        para = footer.paragraphs[0] if idx == 0 else footer.add_paragraph()
        para.alignment = WD_ALIGN_PARAGRAPH.LEFT
        _style_paragraph(para, before_pt=0.0, after_pt=0.0, line=1.4)
        _add_runs(para, txt, bold=bold, color=WORD_COVER_META_RGB, size_pt=spec["meta_size"])


def _add_running_header(section, title: str) -> None:
    """Laufende Kopfzeile: Titel grau + Unterlinie. Sektion wird vom Vorgänger entkoppelt."""
    header = section.header
    header.is_linked_to_previous = False
    para = header.paragraphs[0]
    para.alignment = WD_ALIGN_PARAGRAPH.LEFT
    _style_paragraph(para, before_pt=0.0, after_pt=0.0, line=1.2)
    _add_bottom_border(para, color_hex=WORD_HEADER_LINE_HEX, size=6)
    _add_runs(para, title, bold=False, color=WORD_HEADER_TEXT_RGB, size_pt=WORD_SPEC["header"]["size"])


def build_document(outline: dict) -> bytes:
    doc = Document()
    _apply_a4_section(doc.sections[0])

    normal = doc.styles["Normal"]
    normal.font.name = WORD_FONT_NAME
    normal.font.size = DocxPt(WORD_SPEC["body"]["size"])
    normal.font.color.rgb = WORD_BODY_RGB
    normal_rpr = normal.element.get_or_add_rPr()
    normal_rfonts = normal_rpr.find(qn("w:rFonts"))
    if normal_rfonts is None:
        normal_rfonts = OxmlElement("w:rFonts")
        normal_rpr.append(normal_rfonts)
    for attr in ("w:ascii", "w:hAnsi", "w:cs", "w:eastAsia"):
        normal_rfonts.set(qn(attr), WORD_FONT_NAME)

    title = str(outline.get("title") or "").strip()
    subtitle = str(outline.get("subtitle") or "").strip()
    author = str(outline.get("author") or "").strip()
    date = str(outline.get("date") or "").strip()
    if title:
        # Titelblatt in Sektion 0 (ohne Kopfzeile, Autor/Datum in der Fusszeile), danach eigene
        # Inhaltssektion mit laufender Kopfzeile (Titel grau + Linie) und leerer Fusszeile.
        _add_cover_page(doc, title, subtitle, author, date)
        content_section = doc.add_section(WD_SECTION_START.NEW_PAGE)
        _apply_a4_section(content_section)
        _add_running_header(content_section, title)
        content_section.footer.is_linked_to_previous = False

    blocks = outline.get("blocks") or []
    for block in blocks[:MAX_WORD_BLOCKS]:
        if not isinstance(block, dict):
            continue
        btype = block.get("type")
        if btype == "heading":
            _add_heading_block(doc, block)
        elif btype == "paragraph":
            _add_paragraph_block(doc, block)
        elif btype == "list":
            _add_list_block(doc, block)
        elif btype == "table":
            _add_table_block(doc, block)

    buffer = BytesIO()
    doc.save(buffer)
    return buffer.getvalue()


class RenderRequest(BaseModel):
    html: str


class WordOutlineModel(BaseModel):
    version: int = 1
    fileName: str | None = None
    title: str | None = None
    subtitle: str | None = None
    author: str | None = None
    date: str | None = None
    blocks: list[dict]


class WordRenderRequest(BaseModel):
    outline: WordOutlineModel


app = FastAPI(title="Straton PPTX Renderer")


@app.get("/health")
def health() -> dict:
    return {"status": "ok"}


@app.post("/render")
def render(payload: RenderRequest, x_internal_token: str = Header(default="")) -> Response:
    if not INTERNAL_TOKEN or x_internal_token != INTERNAL_TOKEN:
        raise HTTPException(status_code=401, detail="Ungültiges Internal-Token.")

    slides = parse_slides(payload.html)
    if not slides:
        raise HTTPException(status_code=400, detail="Kein gültiges Folien-HTML gefunden.")
    if len(slides) > MAX_SLIDES:
        raise HTTPException(status_code=400, detail=f"Zu viele Folien (max. {MAX_SLIDES}).")

    pptx_bytes = build_presentation(slides)
    return Response(
        content=pptx_bytes,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


@app.post("/render-docx")
def render_docx(payload: WordRenderRequest, x_internal_token: str = Header(default="")) -> Response:
    if not INTERNAL_TOKEN or x_internal_token != INTERNAL_TOKEN:
        raise HTTPException(status_code=401, detail="Ungültiges Internal-Token.")

    outline = payload.outline.model_dump()
    if not outline.get("blocks"):
        raise HTTPException(status_code=400, detail="Kein gültiger Dokumentinhalt gefunden.")

    docx_bytes = build_document(outline)
    return Response(
        content=docx_bytes,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )
